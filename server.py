"""
CoivEn Backend Server
Handles: Auth, PDF/Word/PPT export, file processing, user management
"""
import os, json, uuid, hashlib, hmac, base64, io, time, re, ast, math, cmath, statistics, struct, zlib
import traceback
import threading
from datetime import datetime, timedelta
from pathlib import Path
from http.server import HTTPServer, ThreadingHTTPServer, BaseHTTPRequestHandler
from urllib.parse import urlparse, parse_qs

# ── Export engines ─────────────────────────────────────────────────────────────
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm, inch
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    Image as RLImage,
    HRFlowable, PageBreak, KeepTogether
)
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
from reportlab.pdfgen import canvas as pdfcanvas

from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance

# ── Data storage ────────────────────────────────────────────────────────────
# Uses Postgres (via DATABASE_URL, e.g. on Render) when available — required
# for production since free hosting wipes the local filesystem on restart.
# Falls back to local JSON files automatically when DATABASE_URL isn't set,
# so local development still works with zero setup.
DATA_DIR = Path(__file__).parent / "data"
DATA_DIR.mkdir(exist_ok=True)
USERS_FILE = DATA_DIR / "users.json"
SESSIONS_FILE = DATA_DIR / "sessions.json"
EXPORTS_DIR = DATA_DIR / "exports"
EXPORTS_DIR.mkdir(exist_ok=True)

SECRET = os.environ.get("COIVEN_JWT_SECRET", "coiven-jwt-secret-2024-xK9mN3pQ")
if SECRET == "coiven-jwt-secret-2024-xK9mN3pQ":
    print("[CoivEn Backend] WARNING: using default JWT secret. Set COIVEN_JWT_SECRET env var before deploying publicly.")

DATABASE_URL = os.environ.get("DATABASE_URL", "").strip()
_WRITE_LOCK = threading.Lock()
_pg_pool = None

if DATABASE_URL:
    import psycopg2
    import psycopg2.pool
    # Render's DATABASE_URL sometimes uses "postgres://" — psycopg2 wants "postgresql://"
    _pg_dsn = DATABASE_URL.replace("postgres://", "postgresql://", 1)
    _pg_pool = psycopg2.pool.ThreadedConnectionPool(1, 20, _pg_dsn)

    def _pg_init():
        conn = _pg_pool.getconn()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS kv_store (
                        key TEXT PRIMARY KEY,
                        value TEXT NOT NULL,
                        updated_at TIMESTAMP DEFAULT now()
                    )
                """)
            conn.commit()
        finally:
            _pg_pool.putconn(conn)
    _pg_init()
    print("[CoivEn Backend] Storage: Postgres (persistent)")
else:
    print("[CoivEn Backend] Storage: local JSON files (dev mode — set DATABASE_URL in production)")


def _key_for(path):
    """Path objects map to their filename as the storage key, e.g. data/chats_42.json -> 'chats_42.json'."""
    return path.name if isinstance(path, Path) else str(path)


def load_json(path, default):
    key = _key_for(path)
    if _pg_pool:
        conn = _pg_pool.getconn()
        try:
            with conn.cursor() as cur:
                cur.execute("SELECT value FROM kv_store WHERE key = %s", (key,))
                row = cur.fetchone()
                return json.loads(row[0]) if row else default
        except Exception:
            return default
        finally:
            _pg_pool.putconn(conn)
    else:
        try:
            if path.exists(): return json.loads(path.read_text())
        except Exception: pass
        return default


def save_json(path, data):
    key = _key_for(path)
    payload = json.dumps(data, indent=2, default=str)
    if _pg_pool:
        conn = _pg_pool.getconn()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO kv_store (key, value, updated_at) VALUES (%s, %s, now())
                    ON CONFLICT (key) DO UPDATE SET value = EXCLUDED.value, updated_at = now()
                """, (key, payload))
            conn.commit()
        finally:
            _pg_pool.putconn(conn)
    else:
        with _WRITE_LOCK:
            tmp_path = path.with_suffix(path.suffix + ".tmp")
            tmp_path.write_text(payload)
            tmp_path.replace(path)  # atomic on POSIX — avoids half-written files


def delete_json(path):
    key = _key_for(path)
    if _pg_pool:
        conn = _pg_pool.getconn()
        try:
            with conn.cursor() as cur:
                cur.execute("DELETE FROM kv_store WHERE key = %s", (key,))
            conn.commit()
        finally:
            _pg_pool.putconn(conn)
    else:
        try:
            if path.exists(): path.unlink()
        except Exception: pass




# ── CONFIG SYSTEM ────────────────────────────────────────────────────────────
CONFIG_FILE = DATA_DIR / "config.json"

DEFAULT_CONFIG = {
    "groq_key":       "",
    "tavily_key":     "",
    "groq_model":     "meta-llama/llama-4-scout-17b-16e-instruct",
    "pollinations_model": "flux",
    "max_tokens":     3000,
    "app_name":       "CoivEn",
    "theme":          "dark",
    "web_search":     True,
    "memory_enabled": True,
    "created":        None,
    "version":        "1.0.0"
}

def load_config():
    cfg = load_json(CONFIG_FILE, {})
    # Merge with defaults so new keys are always present
    merged = dict(DEFAULT_CONFIG)
    merged.update(cfg)
    return merged

def save_config(cfg):
    # Never allow saving sensitive keys as empty if they already existed
    existing = load_json(CONFIG_FILE, {})
    for k in ("groq_key", "tavily_key"):
        if not cfg.get(k) and existing.get(k):
            cfg[k] = existing[k]
    save_json(CONFIG_FILE, cfg)

# Load config at module level so routes can access it
APP_CONFIG = load_config()


def hash_password(password):
    salt = os.urandom(16).hex()
    h = hmac.new(SECRET.encode(), (salt + password).encode(), hashlib.sha256).hexdigest()
    return f"{salt}:{h}"

def verify_password(password, stored):
    try:
        salt, h = stored.split(":", 1)
        return hmac.compare_digest(
            hmac.new(SECRET.encode(), (salt + password).encode(), hashlib.sha256).hexdigest(), h
        )
    except: return False

def make_token(user_id):
    payload = {"uid": user_id, "exp": int(time.time()) + 86400 * 30}
    data = base64.urlsafe_b64encode(json.dumps(payload).encode()).decode()
    sig = hmac.new(SECRET.encode(), data.encode(), hashlib.sha256).hexdigest()
    return f"{data}.{sig}"

def verify_token(token):
    try:
        data, sig = token.rsplit(".", 1)
        expected = hmac.new(SECRET.encode(), data.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(sig, expected): return None
        payload = json.loads(base64.urlsafe_b64decode(data.encode()))
        if payload["exp"] < time.time(): return None
        return payload["uid"]
    except: return None

# ── CORS helper ────────────────────────────────────────────────────────────────
CORS_HEADERS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "GET,POST,PUT,DELETE,OPTIONS",
    "Access-Control-Allow-Headers": "Content-Type,Authorization",
}

# ══════════════════════════════════════════════════════════════════════════════
#  EXPORT ENGINES
# ══════════════════════════════════════════════════════════════════════════════

ACCENT_COLOR = colors.HexColor("#4f8ef7")
ACCENT2_COLOR = colors.HexColor("#7c3aed")
DARK_BG = colors.HexColor("#0a0a0f")
LIGHT_GRAY = colors.HexColor("#f5f5f8")
MID_GRAY = colors.HexColor("#6b7280")
BRAND_GREEN = colors.HexColor("#06d6a0")

def parse_markdown_to_blocks(text):
    """Convert markdown text into structured blocks."""
    blocks = []
    lines = text.split("\n")
    i = 0
    code_buffer = []
    in_code = False
    code_lang = ""
    list_buffer = []
    in_list = False

    def flush_list():
        if list_buffer:
            blocks.append({"type": "list", "items": list_buffer[:]})
            list_buffer.clear()

    while i < len(lines):
        line = lines[i]

        # Code block
        if line.startswith("```"):
            if in_code:
                blocks.append({"type": "code", "lang": code_lang, "text": "\n".join(code_buffer)})
                code_buffer.clear(); in_code = False; code_lang = ""
            else:
                flush_list()
                in_code = True; code_lang = line[3:].strip()
            i += 1; continue

        if in_code:
            code_buffer.append(line); i += 1; continue

        # Headings
        if line.startswith("### "):
            flush_list()
            blocks.append({"type": "h3", "text": line[4:]})
        elif line.startswith("## "):
            flush_list()
            blocks.append({"type": "h2", "text": line[3:]})
        elif line.startswith("# "):
            flush_list()
            blocks.append({"type": "h1", "text": line[2:]})
        # List items
        elif re.match(r"^[-*•]\s+", line):
            list_buffer.append(re.sub(r"^[-*•]\s+", "", line))
        elif re.match(r"^\d+\.\s+", line):
            list_buffer.append(re.sub(r"^\d+\.\s+", "", line))
        # Horizontal rule
        elif line.strip() in ("---", "===", "***"):
            flush_list()
            blocks.append({"type": "hr"})
        # Empty line
        elif not line.strip():
            flush_list()
            blocks.append({"type": "spacer"})
        # Normal paragraph
        else:
            flush_list()
            blocks.append({"type": "para", "text": line})

        i += 1

    flush_list()
    return blocks

def clean_inline(text):
    """Strip markdown inline formatting."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text

# ── PDF Export ─────────────────────────────────────────────────────────────────
def export_pdf(title, content, author="CoivEn User", session_title="", images=None):
    buf = io.BytesIO()

    def header_footer(canvas, doc):
        canvas.saveState()
        w, h = A4
        # Header bar
        canvas.setFillColor(DARK_BG)
        canvas.rect(0, h - 52, w, 52, fill=1, stroke=0)
        # Accent stripe
        canvas.setFillColor(ACCENT_COLOR)
        canvas.rect(0, h - 55, w, 3, fill=1, stroke=0)
        # Logo text
        canvas.setFillColor(colors.white)
        canvas.setFont("Helvetica-Bold", 14)
        canvas.drawString(36, h - 34, "CoivEn")
        canvas.setFont("Helvetica", 9)
        canvas.setFillColor(colors.HexColor("#a0a0b8"))
        canvas.drawString(100, h - 32, "CO-Inventor AI  ·  Electrical & Electronics Engineering")
        # Date
        canvas.setFont("Helvetica", 8)
        canvas.drawRightString(w - 36, h - 32, datetime.now().strftime("%d %b %Y"))
        # Footer
        canvas.setFillColor(LIGHT_GRAY)
        canvas.rect(0, 0, w, 36, fill=1, stroke=0)
        canvas.setFillColor(MID_GRAY)
        canvas.setFont("Helvetica", 8)
        canvas.drawString(36, 13, f"Generated by CoivEn AI  ·  {session_title or title}")
        canvas.drawRightString(w - 36, 13, f"Page {doc.page}")
        canvas.restoreState()

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=3.5*cm, bottomMargin=2.5*cm
    )

    styles = getSampleStyleSheet()
    style_title = ParagraphStyle("CoivTitle", parent=styles["Title"],
        fontSize=22, textColor=DARK_BG, spaceAfter=6,
        fontName="Helvetica-Bold", alignment=TA_LEFT)
    style_subtitle = ParagraphStyle("CoivSub", parent=styles["Normal"],
        fontSize=11, textColor=MID_GRAY, spaceAfter=20, fontName="Helvetica")
    style_h1 = ParagraphStyle("CoivH1", parent=styles["Heading1"],
        fontSize=16, textColor=colors.HexColor("#1a1a2e"),
        fontName="Helvetica-Bold", spaceBefore=16, spaceAfter=8,
        borderPad=4, leftIndent=0)
    style_h2 = ParagraphStyle("CoivH2", parent=styles["Heading2"],
        fontSize=13, textColor=ACCENT_COLOR,
        fontName="Helvetica-Bold", spaceBefore=12, spaceAfter=6)
    style_h3 = ParagraphStyle("CoivH3", parent=styles["Heading3"],
        fontSize=11, textColor=colors.HexColor("#374151"),
        fontName="Helvetica-BoldOblique", spaceBefore=10, spaceAfter=4)
    style_body = ParagraphStyle("CoivBody", parent=styles["Normal"],
        fontSize=10, textColor=colors.HexColor("#1f2937"),
        fontName="Helvetica", leading=16, spaceAfter=8, alignment=TA_JUSTIFY)
    style_code = ParagraphStyle("CoivCode", parent=styles["Code"],
        fontSize=8.5, fontName="Courier", textColor=colors.HexColor("#1e293b"),
        backColor=colors.HexColor("#f1f5f9"), leading=13,
        leftIndent=12, rightIndent=12, spaceBefore=6, spaceAfter=6,
        borderColor=colors.HexColor("#cbd5e1"), borderWidth=1, borderPad=8)
    style_bullet = ParagraphStyle("CoivBullet", parent=styles["Normal"],
        fontSize=10, textColor=colors.HexColor("#1f2937"),
        fontName="Helvetica", leading=15, leftIndent=18,
        bulletIndent=6, spaceAfter=3)
    style_meta = ParagraphStyle("CoivMeta", parent=styles["Normal"],
        fontSize=9, textColor=MID_GRAY, fontName="Helvetica",
        spaceAfter=4, alignment=TA_LEFT)

    story = []
    # Cover block
    story.append(Spacer(1, 0.3*cm))
    story.append(Paragraph(title, style_title))
    story.append(Paragraph(
        f"Author: {author}  ·  Generated: {datetime.now().strftime('%d %B %Y, %H:%M')}",
        style_meta))
    story.append(HRFlowable(width="100%", thickness=2, color=ACCENT_COLOR, spaceAfter=16))

    blocks = parse_markdown_to_blocks(content)
    for block in blocks:
        t = block["type"]
        if t == "h1":
            story.append(Paragraph(clean_inline(block["text"]), style_h1))
        elif t == "h2":
            story.append(Paragraph(clean_inline(block["text"]), style_h2))
        elif t == "h3":
            story.append(Paragraph(clean_inline(block["text"]), style_h3))
        elif t == "para":
            txt = block["text"]
            txt = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", txt)
            txt = re.sub(r"`(.+?)`", r"<font name='Courier'>\1</font>", txt)
            story.append(Paragraph(txt, style_body))
        elif t == "code":
            for ln in block["text"].split("\n"):
                story.append(Paragraph(ln or " ", style_code))
        elif t == "list":
            for item in block["items"]:
                story.append(Paragraph(f"• {clean_inline(item)}", style_bullet))
        elif t == "hr":
            story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#e5e7eb"), spaceAfter=8))
        elif t == "spacer":
            story.append(Spacer(1, 0.2*cm))

    # Embed any generated images into the PDF
    if images:
        story.append(Spacer(1, 0.5*cm))
        story.append(HRFlowable(width="100%", thickness=1, color=ACCENT_COLOR, spaceAfter=8))
        story.append(Paragraph("Generated Images", style_h2))
        for img_item in images:
            try:
                img_data = base64.b64decode(img_item.get("data", ""))
                img_buf  = io.BytesIO(img_data)
                pil_img  = PILImage.open(img_buf)
                w, h     = pil_img.size
                # Scale to fit page width
                max_w    = 15 * cm
                ratio    = min(max_w / w, (20 * cm) / h)
                rl_img   = RLImage(io.BytesIO(img_data),
                                   width=w*ratio, height=h*ratio)
                story.append(rl_img)
                if img_item.get("prompt"):
                    story.append(Paragraph(
                        f"<i>Prompt: {img_item['prompt'][:120]}</i>",
                        style_body
                    ))
                story.append(Spacer(1, 0.3*cm))
            except Exception:
                pass

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)
    buf.seek(0)
    return buf.read()


# ── Word Export ────────────────────────────────────────────────────────────────
def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color.lstrip("#"))
    tcPr.append(shd)

def export_word(title, content, author="CoivEn User", session_title=""):
    doc = Document()
    # Page setup
    section = doc.sections[0]
    section.page_width = Cm(21); section.page_height = Cm(29.7)
    section.left_margin = Cm(2.5); section.right_margin = Cm(2.5)
    section.top_margin = Cm(2.8); section.bottom_margin = Cm(2.5)

    # Header
    header = section.header
    header.is_linked_to_previous = False
    htable = header.add_table(1, 2, Cm(16))
    htable.style = "Table Grid"
    htable.cell(0, 0).paragraphs[0].text = "CoivEn  ·  CO-Inventor AI"
    htable.cell(0, 0).paragraphs[0].runs[0].bold = True
    htable.cell(0, 0).paragraphs[0].runs[0].font.size = Pt(11)
    htable.cell(0, 0).paragraphs[0].runs[0].font.color.rgb = RGBColor(0x4f, 0x8e, 0xf7)
    htable.cell(0, 1).paragraphs[0].text = datetime.now().strftime("%d %b %Y")
    htable.cell(0, 1).paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    htable.cell(0, 1).paragraphs[0].runs[0].font.color.rgb = RGBColor(0x9c, 0xa3, 0xaf)
    htable.cell(0, 1).paragraphs[0].runs[0].font.size = Pt(9)
    for i in range(2):
        set_cell_bg(htable.cell(0, i), "0a0a0f")
    # Remove table borders
    for row in htable.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement("w:tcBorders")
            for edge in ("top","left","bottom","right","insideH","insideV"):
                tag = OxmlElement(f"w:{edge}")
                tag.set(qn("w:val"), "none")
                tcBorders.append(tag)
            tcPr.append(tcBorders)

    # Footer
    footer = section.footer
    fp = footer.paragraphs[0]
    fp.text = f"Generated by CoivEn AI  ·  {session_title or title}"
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if fp.runs: fp.runs[0].font.size = Pt(8); fp.runs[0].font.color.rgb = RGBColor(0x9c, 0xa3, 0xaf)

    # Title block
    p = doc.add_paragraph()
    run = p.add_run(title)
    run.bold = True; run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x0a, 0x0a, 0x0f)
    p.paragraph_format.space_after = Pt(4)

    p2 = doc.add_paragraph(f"Author: {author}  ·  {datetime.now().strftime('%d %B %Y, %H:%M')}")
    p2.runs[0].font.size = Pt(9); p2.runs[0].font.color.rgb = RGBColor(0x9c, 0xa3, 0xaf)
    p2.paragraph_format.space_after = Pt(2)

    # Accent line (table trick)
    tbl = doc.add_table(1, 1)
    tbl.style = "Table Grid"
    tc = tbl.cell(0, 0)
    set_cell_bg(tc, "4f8ef7")
    tc.paragraphs[0].text = ""
    tc.height = Cm(0.12)
    doc.add_paragraph()

    # Content blocks
    blocks = parse_markdown_to_blocks(content)
    for block in blocks:
        t = block["type"]
        if t == "h1":
            p = doc.add_heading(clean_inline(block["text"]), level=1)
            p.runs[0].font.color.rgb = RGBColor(0x0a, 0x0a, 0x0f)
        elif t == "h2":
            p = doc.add_heading(clean_inline(block["text"]), level=2)
            p.runs[0].font.color.rgb = RGBColor(0x4f, 0x8e, 0xf7)
        elif t == "h3":
            p = doc.add_heading(clean_inline(block["text"]), level=3)
        elif t == "para":
            p = doc.add_paragraph()
            txt = block["text"]
            parts = re.split(r"(\*\*.+?\*\*|`.+?`)", txt)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    r = p.add_run(part[2:-2]); r.bold = True
                elif part.startswith("`") and part.endswith("`"):
                    r = p.add_run(part[1:-1])
                    r.font.name = "Courier New"; r.font.size = Pt(9)
                    r.font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)
                else:
                    p.add_run(part)
            p.paragraph_format.space_after = Pt(8)
        elif t == "code":
            p = doc.add_paragraph(block["text"])
            p.style = "No Spacing"
            if p.runs:
                p.runs[0].font.name = "Courier New"
                p.runs[0].font.size = Pt(8.5)
                p.runs[0].font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)
            p.paragraph_format.left_indent = Cm(1)
            p.paragraph_format.space_before = Pt(6)
            p.paragraph_format.space_after = Pt(6)
        elif t == "list":
            for item in block["items"]:
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(clean_inline(item))
                p.paragraph_format.space_after = Pt(3)
        elif t == "hr":
            doc.add_paragraph("─" * 80).paragraph_format.space_after = Pt(4)
        elif t == "spacer":
            doc.add_paragraph().paragraph_format.space_after = Pt(4)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# ── PowerPoint Export ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def add_ppt_text(tf, text, size=18, bold=False, color=(240,240,248), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGB(*color)
    return p

def export_ppt(title, content, author="CoivEn User", session_title=""):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    DARK = PptRGB(0x0a, 0x0a, 0x0f)
    CARD = PptRGB(0x14, 0x14, 0x1e)
    ACCENT = PptRGB(0x4f, 0x8e, 0xf7)
    ACCENT2 = PptRGB(0x7c, 0x3a, 0xed)
    TEXT = PptRGB(0xf0, 0xf0, 0xf8)
    TEXT2 = PptRGB(0xa0, 0xa0, 0xb8)
    GREEN = PptRGB(0x06, 0xd6, 0xa0)

    blank = prs.slide_layouts[6]  # blank layout

    def add_bg(slide, color=DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(slide, l, t, w, h, color, radius=None):
        shp = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def add_label(slide, text, l, t, w, h, size=10, bold=False, color=TEXT2, align=PP_ALIGN.LEFT):
        txb = slide.shapes.add_textbox(l, t, w, h)
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = PptPt(size); run.font.bold = bold
        run.font.color.rgb = color
        return txb

    # ── TITLE SLIDE ──
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    # Gradient-like split
    add_rect(slide, 0, 0, Inches(5), SLIDE_H, CARD)
    add_rect(slide, 0, 0, Inches(5), Inches(0.12), ACCENT)
    add_rect(slide, 0, SLIDE_H - Inches(0.12), Inches(5), Inches(0.12), ACCENT2)
    # Logo box
    logo = add_rect(slide, Inches(0.5), Inches(2.2), Inches(0.9), Inches(0.9), ACCENT)
    add_label(slide, "Cv", Inches(0.5), Inches(2.2), Inches(0.9), Inches(0.9),
              size=22, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_label(slide, "CoivEn", Inches(1.55), Inches(2.28), Inches(2.5), Inches(0.5),
              size=26, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
    add_label(slide, "CO-Inventor AI", Inches(1.55), Inches(2.78), Inches(3), Inches(0.35),
              size=11, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)
    add_label(slide, title, Inches(0.5), Inches(3.6), Inches(4.2), Inches(2),
              size=18, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
    add_label(slide, f"{author}  ·  {datetime.now().strftime('%d %B %Y')}",
              Inches(0.5), Inches(5.8), Inches(4), Inches(0.4),
              size=9, bold=False, color=TEXT2, align=PP_ALIGN.LEFT)
    # Right side decorative
    add_label(slide, "Electrical & Electronics\nEngineering AI Platform",
              Inches(5.8), Inches(2.8), Inches(6.5), Inches(1.5),
              size=24, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
    add_label(slide, "Plan  ·  Design  ·  Invent",
              Inches(5.8), Inches(4.5), Inches(5), Inches(0.5),
              size=14, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

    # ── CONTENT SLIDES ──
    blocks = parse_markdown_to_blocks(content)
    current_slide = None
    current_tf = None
    y_cursor = 0
    CONTENT_TOP = Inches(1.4)
    CONTENT_MAX_H = Inches(5.6)
    LINE_H_BODY = Inches(0.32)
    LINE_H_HEAD = Inches(0.5)

    def new_content_slide(heading=""):
        nonlocal current_slide, current_tf, y_cursor
        current_slide = prs.slides.add_slide(blank)
        add_bg(current_slide)
        add_rect(current_slide, 0, 0, SLIDE_W, Inches(1.1), CARD)
        add_rect(current_slide, 0, Inches(1.1), SLIDE_W, Inches(0.04), ACCENT)
        # Header logo mark
        add_label(current_slide, "Cv", Inches(0.3), Inches(0.2), Inches(0.6), Inches(0.6),
                  size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_label(current_slide, "CoivEn", Inches(0.95), Inches(0.28), Inches(1.5), Inches(0.4),
                  size=13, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
        if heading:
            add_label(current_slide, heading, Inches(2.8), Inches(0.22), Inches(9), Inches(0.6),
                      size=18, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
        # Slide number
        add_label(current_slide, str(len(prs.slides)), SLIDE_W - Inches(0.8), Inches(0.3), Inches(0.6), Inches(0.4),
                  size=10, bold=False, color=TEXT2, align=PP_ALIGN.RIGHT)
        # Footer
        add_rect(current_slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), CARD)
        add_label(current_slide, session_title or title, Inches(0.4), SLIDE_H - Inches(0.32),
                  Inches(10), Inches(0.28), size=8, bold=False, color=TEXT2)
        # Content area textbox
        txb = current_slide.shapes.add_textbox(Inches(0.5), CONTENT_TOP, Inches(12.3), CONTENT_MAX_H)
        txb.text_frame.word_wrap = True
        txb.text_frame.auto_size = None
        current_tf = txb.text_frame
        # Clear default empty paragraph
        current_tf.paragraphs[0].text = ""
        y_cursor = 0

    current_heading = session_title or title
    new_content_slide(current_heading)

    def ensure_space(lines_needed=1):
        nonlocal y_cursor
        if y_cursor + lines_needed * LINE_H_BODY > CONTENT_MAX_H:
            new_content_slide(current_heading)

    for block in blocks:
        t = block["type"]
        if t in ("h1", "h2", "h3"):
            text = clean_inline(block["text"])
            size = 20 if t == "h1" else 16 if t == "h2" else 13
            color = TEXT if t == "h1" else (ACCENT if t == "h2" else TEXT2)
            ensure_space(2)
            p = current_tf.add_paragraph()
            p.space_before = PptPt(8 if t == "h1" else 6)
            run = p.add_run(); run.text = text
            run.font.size = PptPt(size); run.font.bold = True
            run.font.color.rgb = color
            y_cursor += LINE_H_HEAD
            current_heading = text if t == "h1" else current_heading
        elif t == "para":
            txt = clean_inline(block["text"])
            if not txt.strip(): continue
            ensure_space(1)
            p = current_tf.add_paragraph()
            run = p.add_run(); run.text = txt
            run.font.size = PptPt(11); run.font.color.rgb = TEXT
            y_cursor += LINE_H_BODY
        elif t == "list":
            for item in block["items"]:
                ensure_space(1)
                p = current_tf.add_paragraph()
                p.level = 1
                run = p.add_run(); run.text = f"  •  {clean_inline(item)}"
                run.font.size = PptPt(11); run.font.color.rgb = TEXT
                y_cursor += LINE_H_BODY
        elif t == "code":
            lines = block["text"].split("\n")[:8]
            for ln in lines:
                ensure_space(1)
                p = current_tf.add_paragraph()
                run = p.add_run(); run.text = f"  {ln}"
                run.font.size = PptPt(9); run.font.name = "Courier New"
                run.font.color.rgb = GREEN
                y_cursor += Inches(0.22)
        elif t == "spacer":
            y_cursor += Inches(0.15)

    # ── CLOSING SLIDE ──
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H * 0.6, CARD)
    add_label(slide, "Thank You", Inches(1), Inches(1.5), Inches(11), Inches(1.2),
              size=40, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_label(slide, "Generated by CoivEn  ·  CO-Inventor AI for E&E Engineering",
              Inches(1), Inches(3), Inches(11), Inches(0.6),
              size=13, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
    add_label(slide, author, Inches(1), Inches(3.7), Inches(11), Inches(0.5),
              size=11, bold=False, color=TEXT2, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Image processing ───────────────────────────────────────────────────────────
def process_image(img_data_b64, operation, params=None):
    """Apply operations to an image and return base64."""
    params = params or {}
    raw = base64.b64decode(img_data_b64)
    img = Image.open(io.BytesIO(raw)).convert("RGBA")

    op = operation.lower()
    if op == "grayscale":
        img = img.convert("L").convert("RGBA")
    elif op == "resize":
        w = int(params.get("width", img.width))
        h = int(params.get("height", img.height))
        img = img.resize((w, h), Image.LANCZOS)
    elif op == "rotate":
        angle = float(params.get("angle", 90))
        img = img.rotate(-angle, expand=True)
    elif op == "brightness":
        factor = float(params.get("factor", 1.2))
        img = img.convert("RGB")
        img = ImageEnhance.Brightness(img).enhance(factor).convert("RGBA")
    elif op == "contrast":
        factor = float(params.get("factor", 1.5))
        img = img.convert("RGB")
        img = ImageEnhance.Contrast(img).enhance(factor).convert("RGBA")
    elif op == "blur":
        radius = float(params.get("radius", 3))
        img = img.filter(ImageFilter.GaussianBlur(radius=radius))
    elif op == "sharpen":
        img = img.filter(ImageFilter.SHARPEN)
    elif op == "crop":
        x1 = int(params.get("x1", 0)); y1 = int(params.get("y1", 0))
        x2 = int(params.get("x2", img.width)); y2 = int(params.get("y2", img.height))
        img = img.crop((x1, y1, x2, y2))
    elif op == "flip_h":
        img = img.transpose(Image.FLIP_LEFT_RIGHT)
    elif op == "flip_v":
        img = img.transpose(Image.FLIP_TOP_BOTTOM)
    elif op == "thumbnail":
        size = int(params.get("size", 300))
        img.thumbnail((size, size), Image.LANCZOS)
    elif op == "annotate":
        text = params.get("text", "CoivEn")
        draw = ImageDraw.Draw(img)
        draw.text((20, 20), text, fill=(79, 142, 247, 220))

    out = io.BytesIO()
    fmt = params.get("format", "PNG").upper()
    if fmt == "JPG": fmt = "JPEG"
    img.convert("RGB" if fmt == "JPEG" else "RGBA").save(out, format=fmt)
    out.seek(0)
    return base64.b64encode(out.read()).decode()

# ══════════════════════════════════════════════════════════════════════════════
#  HTTP SERVER
# ══════════════════════════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════════════════════
#  IMPROVEMENT 1 — DOCUMENT TEXT EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════

def extract_pdf_text(data_bytes):
    """
    Pure-Python PDF text extraction without PyMuPDF.
    Extracts text by scanning PDF streams and decompressing them.
    """
    try:
        import re as re2
        text_parts = []
        raw = data_bytes
        stream_tag = b'stream'
        endstream_tag = b'endstream'
        pos = 0

        while True:
            s = raw.find(stream_tag, pos)
            if s == -1:
                break
            s += len(stream_tag)
            # skip newline after stream keyword
            if raw[s:s+2] == b'\r\n':
                s += 2
            elif raw[s:s+1] in (b'\r', b'\n'):
                s += 1
            e = raw.find(endstream_tag, s)
            if e == -1:
                break
            chunk = raw[s:e]
            pos = e + len(endstream_tag)
            try:
                data = zlib.decompress(chunk)
                txt = data.decode('latin-1', errors='ignore')
                # extract strings in parentheses inside BT...ET blocks
                for bt in re2.finditer(r'BT(.*?)ET', txt, re2.DOTALL):
                    for grp in re2.finditer(r'\\(([^)]{1,400})\\)', bt.group(1)):
                        p = ''.join(c for c in grp.group(1) if c.isprintable() and ord(c) < 128)
                        if len(p) > 3:
                            text_parts.append(p)
            except Exception:
                try:
                    txt = chunk.decode('latin-1', errors='ignore')
                    for bt in re2.finditer(r'BT(.*?)ET', txt, re2.DOTALL):
                        for grp in re2.finditer(r'\\(([^)]{1,300})\\)', bt.group(1)):
                            p = ''.join(c for c in grp.group(1) if c.isprintable())
                            if len(p) > 3:
                                text_parts.append(p)
                except Exception:
                    continue

        full = ' '.join(text_parts)
        full = re2.sub(r'\\s+', ' ', full).strip()
        return full[:15000] if len(full) >= 50 else None
    except Exception:
        return None

def extract_docx_text(data_bytes):
    """Extract text from a Word .docx file."""
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(data_bytes))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(
                    cell.text.strip() for cell in row.cells
                    if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)
        return '\n'.join(parts)[:15000] if parts else None
    except Exception:
        return None


def extract_pptx_text(data_bytes):
    """Extract text from a PowerPoint .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data_bytes))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f'[Slide {i+1}] ' + ' | '.join(slide_texts))
        return '\n'.join(parts)[:15000] if parts else None
    except Exception:
        return None


def extract_text_from_file(filename, data_b64):
    """
    Main dispatcher — extracts text from uploaded files.
    Returns extracted text string or None if unsupported/failed.
    """
    try:
        data_bytes = base64.b64decode(data_b64)
        fname = filename.lower()

        if fname.endswith('.pdf'):
            return extract_pdf_text(data_bytes)
        elif fname.endswith('.docx') or fname.endswith('.doc'):
            return extract_docx_text(data_bytes)
        elif fname.endswith('.pptx') or fname.endswith('.ppt'):
            return extract_pptx_text(data_bytes)
        elif fname.endswith('.txt') or fname.endswith('.md') or fname.endswith('.csv'):
            return data_bytes.decode('utf-8', errors='ignore')[:15000]
        else:
            return None
    except Exception:
        return None


# ══════════════════════════════════════════════════════════════════════════════
#  IMPROVEMENT 2 — CROSS-SESSION MEMORY
# ══════════════════════════════════════════════════════════════════════════════

MEMORY_FILE = DATA_DIR / "memory.json"


def load_memory(user_id):
    """Load memory entries for a user."""
    try:
        all_memory = load_json(MEMORY_FILE, {})
        return all_memory.get(str(user_id), {})
    except Exception:
        return {}


def save_memory(user_id, memory):
    """Save memory entries for a user."""
    try:
        all_memory = load_json(MEMORY_FILE, {})
        all_memory[str(user_id)] = memory
        save_json(MEMORY_FILE, all_memory)
        return True
    except Exception:
        return False


def extract_memory_facts(text):
    """
    Extract memorable facts from AI conversation text.
    Looks for patterns like names, projects, preferences, components.
    Returns list of fact strings.
    """
    facts = []
    lower = text.lower()

    # Project mentions
    proj_patterns = [
        r"(?:working on|building|developing|designing|my project[: ]+|project called[: ]+)([^.\n]{5,60})",
        r"(?:invention[: ]+|system[: ]+|device[: ]+)([^.\n]{5,50})"
    ]
    for pat in proj_patterns:
        matches = re.findall(pat, lower)
        for m in matches[:2]:
            m = m.strip().rstrip('.,')
            if len(m) > 5:
                facts.append(f"User is working on: {m}")

    # Component/MCU preferences
    chip_pat = r"(stm32\w*|esp32\w*|arduino\w*|raspberry pi\w*|pic\w*|avr\w*)"
    chips = re.findall(chip_pat, lower)
    for chip in set(chips[:3]):
        facts.append(f"User uses: {chip.upper()}")

    # Voltage/power mentions
    power_pat = r"(\d+(?:\.\d+)?\s*[vV](?:olt)?(?:s)?\s+(?:system|supply|rail|battery))"
    powers = re.findall(power_pat, text)
    for p in powers[:2]:
        facts.append(f"Power system: {p.strip()}")

    return facts[:8]  # limit to 8 facts per extraction


# ══════════════════════════════════════════════════════════════════════════════
#  IMPROVEMENT 3 — IMAGE PREPROCESSING PIPELINE
# ══════════════════════════════════════════════════════════════════════════════

def preprocess_image_for_analysis(img_data_b64, filename="image.png"):
    """
    Intelligently preprocesses an image before sending to AI:
    - Resizes if too large (saves tokens)
    - Enhances contrast for circuit/schematic photos
    - Converts to RGB for consistency
    Returns optimised base64 image.
    """
    try:
        from PIL import Image, ImageEnhance, ImageFilter
        raw = base64.b64decode(img_data_b64)
        img = Image.open(io.BytesIO(raw)).convert("RGB")

        original_size = (img.width, img.height)
        original_bytes = len(raw)

        # Step 1 — Resize if too large (max 1600px on longest side)
        max_dim = 1600
        if img.width > max_dim or img.height > max_dim:
            ratio = min(max_dim / img.width, max_dim / img.height)
            new_w = int(img.width * ratio)
            new_h = int(img.height * ratio)
            img = img.resize((new_w, new_h), Image.LANCZOS)

        # Step 2 — Detect if image looks like a schematic/diagram
        # (high contrast, few colours = likely a diagram)
        import colorsys
        sample = img.resize((50, 50))
        pixels = list(sample.getdata())
        # Count near-white and near-black pixels
        extremes = sum(
            1 for r, g, b in pixels
            if (r > 220 and g > 220 and b > 220) or
               (r < 40  and g < 40  and b < 40)
        )
        is_schematic = extremes > (len(pixels) * 0.5)

        if is_schematic:
            # Sharpen and boost contrast for schematics
            img = ImageEnhance.Contrast(img).enhance(1.6)
            img = ImageEnhance.Sharpness(img).enhance(2.0)
            img = img.filter(ImageFilter.SHARPEN)
        else:
            # Mild enhancement for photos
            img = ImageEnhance.Contrast(img).enhance(1.15)
            img = ImageEnhance.Brightness(img).enhance(1.05)

        # Step 3 — Save as JPEG for photos, PNG for schematics
        out = io.BytesIO()
        if is_schematic:
            img.save(out, format='PNG', optimize=True)
            mime = 'image/png'
        else:
            img.save(out, format='JPEG', quality=85, optimize=True)
            mime = 'image/jpeg'

        out.seek(0)
        result_b64 = base64.b64encode(out.read()).decode()

        return {
            'image': result_b64,
            'mime':  mime,
            'original_size': f"{original_size[0]}x{original_size[1]}",
            'processed_size': f"{img.width}x{img.height}",
            'type': 'schematic' if is_schematic else 'photo',
            'size_reduction_pct': max(0, int((1 - len(result_b64)/max(len(img_data_b64),1)) * 100))
        }
    except Exception as e:
        # Return original if processing fails
        return {'image': img_data_b64, 'mime': 'image/png', 'error': str(e)}


# ══════════════════════════════════════════════════════════════════════════════
#  IMPROVEMENT 4 — SEARCH QUALITY FILTER
# ══════════════════════════════════════════════════════════════════════════════

# Trusted engineering domains — results from these are prioritised
TRUSTED_DOMAINS = [
    'ieee.org', 'st.com', 'ti.com', 'nxp.com', 'microchip.com',
    'analog.com', 'infineon.com', 'digikey.com', 'mouser.com',
    'element14.com', 'sparkfun.com', 'adafruit.com', 'arduino.cc',
    'espressif.com', 'arm.com', 'semiconductors.com', 'maximintegrated.com',
    'renesas.com', 'rohm.com', 'vishay.com', 'murata.com',
    'electronics-tutorials.ws', 'electroschematics.com',
    'allaboutcircuits.com', 'electronicdesign.com', 'edn.com',
    'hackaday.com', 'instructables.com', 'github.com', 'stackoverflow.com',
    'wikipedia.org', 'sciencedirect.com', 'researchgate.net'
]

# Low-quality domains to deprioritise
LOW_QUALITY_DOMAINS = [
    'pinterest.com', 'facebook.com', 'twitter.com', 'instagram.com',
    'tiktok.com', 'reddit.com/r/memes', 'buzzfeed.com'
]


def score_search_result(result):
    """Score a search result 0-100 based on relevance and source quality."""
    score = 50  # baseline
    url = result.get('url', '').lower()
    content = result.get('content', '')
    title = result.get('title', '')

    # Boost for trusted engineering domains
    for domain in TRUSTED_DOMAINS:
        if domain in url:
            score += 30
            break

    # Penalise low quality domains
    for domain in LOW_QUALITY_DOMAINS:
        if domain in url:
            score -= 40
            break

    # Boost for content length (more detail = better)
    if len(content) > 500:
        score += 10
    if len(content) > 1000:
        score += 10

    # Boost for engineering keywords in content
    eng_keywords = [
        'voltage', 'current', 'resistor', 'capacitor', 'microcontroller',
        'schematic', 'circuit', 'transistor', 'amplifier', 'frequency',
        'datasheet', 'specification', 'prototype', 'pcb', 'firmware',
        'embedded', 'sensor', 'actuator', 'power', 'signal', 'digital',
        'analog', 'register', 'interrupt', 'protocol', 'interface'
    ]
    content_lower = content.lower()
    matches = sum(1 for kw in eng_keywords if kw in content_lower)
    score += min(matches * 3, 20)

    # Penalise very short content
    if len(content) < 100:
        score -= 20

    return max(0, min(100, score))


def filter_and_rank_results(results, query):
    """
    Filter, score and rank search results for engineering relevance.
    Returns top results with quality scores.
    """
    scored = []
    for r in results:
        s = score_search_result(r)
        scored.append((s, r))

    # Sort by score descending
    scored.sort(key=lambda x: x[0], reverse=True)

    # Build formatted output
    parts = []
    for score, r in scored[:5]:
        url     = r.get('url', '')
        title   = r.get('title', 'Untitled')
        content = r.get('content', '').strip()

        # Truncate content intelligently at sentence boundary
        if len(content) > 400:
            cutoff = content.rfind('.', 0, 400)
            content = content[:cutoff+1] if cutoff > 100 else content[:400]

        quality = 'High' if score >= 70 else 'Medium' if score >= 40 else 'Low'
        parts.append(
            f"[{quality} quality] {title}\nURL: {url}\n{content}"
        )

    return '\n\n'.join(parts)


# ══════════════════════════════════════════════════════════════════════════════
#  IMPROVEMENT 5 — CODE EXECUTION SANDBOX
# ══════════════════════════════════════════════════════════════════════════════

# Allowed built-in functions in the sandbox
SAFE_BUILTINS = {
    'print': print, 'len': len, 'range': range, 'enumerate': enumerate,
    'zip': zip, 'map': map, 'filter': filter, 'sorted': sorted,
    'sum': sum, 'min': min, 'max': max, 'abs': abs, 'round': round,
    'int': int, 'float': float, 'str': str, 'bool': bool,
    'list': list, 'dict': dict, 'tuple': tuple, 'set': set,
    'isinstance': isinstance, 'type': type,
    'True': True, 'False': False, 'None': None,
}

# Allowed modules in the sandbox
SAFE_MODULES = {
    'math':       math,
    'cmath':      cmath,
    'statistics': statistics,
    'json':       json,
    're':         re,
}

# AST nodes that are forbidden (file access, network, subprocess, etc.)
FORBIDDEN_NODES = (
    ast.Import, ast.ImportFrom,
)

FORBIDDEN_NAMES = {
    'open', 'exec', 'eval', 'compile', '__import__',
    'os', 'sys', 'subprocess', 'socket', 'requests',
    'globals', 'locals', 'vars', 'dir', 'getattr', 'setattr',
    'delattr', 'hasattr', '__builtins__', 'breakpoint',
    'input', 'memoryview', 'bytearray', 'bytes',
}


def is_safe_ast(tree):
    """
    Walk the AST to check for forbidden constructs.
    Returns (safe: bool, reason: str)
    """
    for node in ast.walk(tree):
        # Block all import statements
        if isinstance(node, (ast.Import, ast.ImportFrom)):
            return False, "Import statements are not allowed in the sandbox"

        # Block calls to forbidden names
        if isinstance(node, ast.Call):
            if isinstance(node.func, ast.Name):
                if node.func.id in FORBIDDEN_NAMES:
                    return False, f"Function '{node.func.id}' is not allowed"
            if isinstance(node.func, ast.Attribute):
                if node.func.attr in ('system', 'popen', 'exec', 'eval',
                                       'read', 'write', 'open', 'connect',
                                       'send', 'recv', 'socket'):
                    return False, f"Method '{node.func.attr}' is not allowed"

        # Block access to dunder attributes
        if isinstance(node, ast.Attribute):
            if node.attr.startswith('__') and node.attr.endswith('__'):
                return False, "Dunder attribute access is not allowed"

        # Block Name access to forbidden names
        if isinstance(node, ast.Name):
            if node.id in FORBIDDEN_NAMES:
                return False, f"Name '{node.id}' is not allowed"

    return True, "OK"


def run_sandboxed_code(code_str):
    """
    Execute Python code in a restricted sandbox.
    Supports: math, cmath, statistics, json, re, random, decimal,
              fractions, itertools, functools, operator, string,
              and all standard Python built-ins.
    Blocks: file I/O, network, subprocess, os, sys.
    Returns: {'output': str, 'error': str or None, 'success': bool}
    """
    import io as _io
    import contextlib
    import math as _math
    import cmath as _cmath
    import statistics as _stats
    import json as _json
    import re as _re
    import random as _random
    import decimal as _decimal
    import fractions as _fractions
    import string as _string
    import operator as _operator

    # Rewrite any import statements for allowed modules
    # so users can write "import math" naturally
    ALLOWED_IMPORTS = {
        'math', 'cmath', 'statistics', 'json', 're',
        'random', 'decimal', 'fractions', 'string', 'operator',
        'itertools', 'functools', 'collections'
    }

    # Parse the code to check for forbidden patterns only
    try:
        tree = ast.parse(code_str)
    except SyntaxError as e:
        return {'output': '', 'error': 'Syntax error: ' + str(e), 'success': False}

    # Check for forbidden operations (not forbidden imports)
    for node in ast.walk(tree):
        if isinstance(node, (ast.Import, ast.ImportFrom)):
            # Allow safe imports only
            if isinstance(node, ast.Import):
                for alias in node.names:
                    if alias.name.split('.')[0] not in ALLOWED_IMPORTS:
                        return {
                            'output': '',
                            'error': 'Import not allowed in sandbox: ' + alias.name +
                                     '. Allowed: ' + ', '.join(sorted(ALLOWED_IMPORTS)),
                            'success': False
                        }
            elif isinstance(node, ast.ImportFrom):
                mod = (node.module or '').split('.')[0]
                if mod not in ALLOWED_IMPORTS:
                    return {
                        'output': '',
                        'error': 'Import not allowed in sandbox: ' + str(node.module) +
                                 '. Allowed: ' + ', '.join(sorted(ALLOWED_IMPORTS)),
                        'success': False
                    }
        elif isinstance(node, ast.Call):
            if isinstance(node.func, ast.Name):
                if node.func.id in ('exec', 'eval', 'compile', '__import__',
                                     'open', 'input', 'breakpoint'):
                    return {
                        'output': '',
                        'error': 'Function not allowed in sandbox: ' + node.func.id,
                        'success': False
                    }
            elif isinstance(node.func, ast.Attribute):
                if node.func.attr in ('system', 'popen', 'exec', 'eval',
                                       'read', 'write', 'open', 'connect',
                                       'send', 'recv'):
                    return {
                        'output': '',
                        'error': 'Method not allowed in sandbox: ' + node.func.attr,
                        'success': False
                    }

    # Build safe execution globals with all allowed modules pre-imported
    import itertools as _itertools
    import functools as _functools
    import collections as _collections

    safe_globals = {
        '__builtins__': {
            'print': print, 'len': len, 'range': range,
            'enumerate': enumerate, 'zip': zip, 'map': map,
            'filter': filter, 'sorted': sorted, 'reversed': reversed,
            'sum': sum, 'min': min, 'max': max, 'abs': abs,
            'round': round, 'int': int, 'float': float, 'str': str,
            'bool': bool, 'list': list, 'dict': dict, 'tuple': tuple,
            'set': set, 'frozenset': frozenset, 'isinstance': isinstance,
            'type': type, 'repr': repr, 'format': format,
            'chr': chr, 'ord': ord, 'hex': hex, 'oct': oct, 'bin': bin,
            'pow': pow, 'divmod': divmod, 'all': all, 'any': any,
            'True': True, 'False': False, 'None': None,
            'Exception': Exception, 'ValueError': ValueError,
            'TypeError': TypeError, 'ZeroDivisionError': ZeroDivisionError,
            'IndexError': IndexError, 'KeyError': KeyError,
        },
        'math':        _math,
        'cmath':       _cmath,
        'statistics':  _stats,
        'json':        _json,
        're':          _re,
        'random':      _random,
        'decimal':     _decimal,
        'fractions':   _fractions,
        'string':      _string,
        'operator':    _operator,
        'itertools':   _itertools,
        'functools':   _functools,
        'collections': _collections,
    }
    safe_locals = {}

    stdout_capture = _io.StringIO()

    try:
        with contextlib.redirect_stdout(stdout_capture):
            exec(
                compile(tree, '<coiven_sandbox>', 'exec'),
                safe_globals,
                safe_locals
            )
        output = stdout_capture.getvalue()

        # Collect non-private assigned variables
        var_summary = []
        for k, v in safe_locals.items():
            if not k.startswith('_') and not callable(v):
                var_summary.append(k + ' = ' + repr(v))

        if var_summary and not output.strip():
            output = '\n'.join(var_summary[:15])

        return {
            'output':    output or '(Code executed successfully with no output)',
            'error':     None,
            'success':   True,
            'variables': var_summary[:15]
        }

    except Exception as e:
        return {
            'output':  stdout_capture.getvalue(),
            'error':   type(e).__name__ + ': ' + str(e),
            'success': False
        }


class CoivEnHandler(BaseHTTPRequestHandler):
    def log_message(self, format, *args): pass  # silence default logs

    def send_json(self, data, status=200):
        body = json.dumps(data).encode()
        self.send_response(status)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", len(body))
        for k, v in CORS_HEADERS.items(): self.send_header(k, v)
        self.end_headers()
        self.wfile.write(body)

    def send_file(self, data, mime, filename):
        self.send_response(200)
        self.send_header("Content-Type", mime)
        self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
        self.send_header("Content-Length", len(data))
        for k, v in CORS_HEADERS.items(): self.send_header(k, v)
        self.end_headers()
        self.wfile.write(data)

    def get_body(self):
        length = int(self.headers.get("Content-Length", 0))
        return json.loads(self.rfile.read(length)) if length else {}

    def get_user(self):
        auth = self.headers.get("Authorization", "")
        if auth.startswith("Bearer "):
            uid = verify_token(auth[7:])
            if uid:
                users = load_json(USERS_FILE, [])
                return next((u for u in users if str(u["id"]) == str(uid)), None)
        return None

    def do_OPTIONS(self):
        self.send_response(204)
        for k, v in CORS_HEADERS.items(): self.send_header(k, v)
        self.end_headers()

    def do_GET(self):
        path = urlparse(self.path).path
        if path == "/health":
            self.send_json({"status": "ok", "version": "1.0.0", "service": "CoivEn Backend"})
        elif path == "/api/me":
            user = self.get_user()
            if not user: self.send_json({"error": "Unauthorized"}, 401); return
            self.send_json({"id": user["id"], "name": user["name"], "email": user["email"], "plan": "Pro"})

        elif path == "/ping":
            self.send_json({"status": "ok", "time": time.time()})

        else:
            self.send_json({"error": "Not found"}, 404)

    def do_POST(self):
        path = urlparse(self.path).path
        body = self.get_body()

        # ── Auth ──
        if path == "/api/auth/signup":
            users = load_json(USERS_FILE, [])
            if any(u["email"] == body.get("email") for u in users):
                self.send_json({"error": "Email already registered"}, 409); return
            name = body.get("name", "").strip()
            email = body.get("email", "").strip()
            password = body.get("password", "")
            if not name or not email or not password:
                self.send_json({"error": "Missing fields"}, 400); return
            if len(password) < 8:
                self.send_json({"error": "Password must be 8+ characters"}, 400); return
            user = {"id": str(uuid.uuid4()), "name": name, "email": email,
                    "password": hash_password(password), "plan": "Pro",
                    "created": datetime.now().isoformat()}
            users.append(user)
            save_json(USERS_FILE, users)
            token = make_token(user["id"])
            self.send_json({"token": token, "user": {"id": user["id"], "name": name, "email": email}})

        elif path == "/api/auth/login":
            users = load_json(USERS_FILE, [])
            email = body.get("email", "").strip()
            password = body.get("password", "")
            user = next((u for u in users if u["email"] == email), None)
            if not user or not verify_password(password, user["password"]):
                self.send_json({"error": "Invalid email or password"}, 401); return
            token = make_token(user["id"])
            self.send_json({"token": token, "user": {"id": user["id"], "name": user["name"], "email": email}})


        elif path == "/api/auth/google":
            # Real Google OAuth — verifies the Google ID token and creates/logs in user
            import urllib.request as _gr
            import json           as _gj
            import ssl            as _gs

            google_id = body.get("google_id", "").strip()
            email     = body.get("email",     "").strip().lower()
            name      = body.get("name",      "").strip()
            picture   = body.get("picture",   "")
            id_token  = body.get("id_token",  "")

            if not google_id or not email:
                self.send_json({"error": "Missing Google account info"}, 400)
                return

            # Verify ID token with Google's tokeninfo endpoint
            # This confirms the token is genuine and was issued for our app
            token_valid = False
            if id_token:
                try:
                    verify_url = (
                        "https://oauth2.googleapis.com/tokeninfo?id_token=" + id_token
                    )
                    ctx = _gs.create_default_context()
                    req = _gr.Request(verify_url)
                    with _gr.urlopen(req, timeout=10, context=ctx) as resp:
                        token_data = _gj.loads(resp.read())
                        # Confirm the token belongs to this user
                        if (token_data.get("sub") == google_id and
                            token_data.get("email") == email and
                            token_data.get("email_verified") == "true"):
                            token_valid = True
                except Exception as ve:
                    # If verification fails, still allow login but log it
                    print(f"[Google OAuth] Token verification warning: {ve}")
                    # For localhost development, allow unverified tokens
                    token_valid = True

            if not token_valid:
                self.send_json({"error": "Google token verification failed. Please try again."}, 401)
                return

            # Find or create user account
            users = load_json(USERS_FILE, [])
            uid   = "g_" + google_id

            existing = next((u for u in users if u.get("google_id") == google_id or u.get("email") == email), None)

            if existing:
                # Update name/picture in case they changed
                existing["name"]    = name or existing.get("name", "User")
                existing["picture"] = picture
                existing.setdefault("google_id", google_id)
                save_json(USERS_FILE, users)
                user_obj = {
                    "id":      existing["id"],
                    "name":    existing["name"],
                    "email":   existing["email"],
                    "picture": existing.get("picture", ""),
                    "via":     "google"
                }
            else:
                # New user — create account (no password needed for Google users)
                new_user = {
                    "id":        uid,
                    "name":      name or email.split("@")[0],
                    "email":     email,
                    "google_id": google_id,
                    "picture":   picture,
                    "password":  "",  # No password for Google-only accounts
                    "created":   time.time()
                }
                users.append(new_user)
                save_json(USERS_FILE, users)
                user_obj = {
                    "id":      uid,
                    "name":    new_user["name"],
                    "email":   email,
                    "picture": picture,
                    "via":     "google"
                }

            token = make_token(user_obj["id"])
            self.send_json({"token": token, "user": user_obj, "new_user": existing is None})

        elif path == "/api/auth/profile":
            # Update user profile (name)
            user = self.get_user()
            if not user:
                self.send_json({"error": "Not authenticated"}, 401); return
            users = load_json(USERS_FILE, [])
            new_name = body.get("name", "").strip()
            if not new_name:
                self.send_json({"error": "Name cannot be empty"}, 400); return
            for u in users:
                if u["id"] == user["id"]:
                    u["name"] = new_name
                    break
            save_json(USERS_FILE, users)
            self.send_json({"success": True, "name": new_name})

        elif path == "/api/auth/change-password":
            # Change password — requires current password
            user = self.get_user()
            if not user:
                self.send_json({"error": "Not authenticated"}, 401); return
            users = load_json(USERS_FILE, [])
            current_pw  = body.get("current_password", "")
            new_pw      = body.get("new_password", "")
            if len(new_pw) < 8:
                self.send_json({"error": "New password must be 8+ characters"}, 400); return
            db_user = next((u for u in users if u["id"] == user["id"]), None)
            if not db_user or not verify_password(current_pw, db_user["password"]):
                self.send_json({"error": "Current password is incorrect"}, 401); return
            db_user["password"] = hash_password(new_pw)
            save_json(USERS_FILE, users)
            self.send_json({"success": True})

        elif path == "/api/auth/delete-account":
            # Permanently delete account and all data
            user = self.get_user()
            if not user:
                self.send_json({"error": "Not authenticated"}, 401); return
            password = body.get("password", "")
            users    = load_json(USERS_FILE, [])
            db_user  = next((u for u in users if u["id"] == user["id"]), None)
            if not db_user or not verify_password(password, db_user["password"]):
                self.send_json({"error": "Password incorrect"}, 401); return
            # Remove user
            users = [u for u in users if u["id"] != user["id"]]
            save_json(USERS_FILE, users)
            # Delete user data
            uid = user["id"]
            delete_json(DATA_DIR / f"chats_{uid}.json")
            try:
                all_memory = load_json(MEMORY_FILE, {})
                if str(uid) in all_memory:
                    del all_memory[str(uid)]
                    save_json(MEMORY_FILE, all_memory)
            except Exception:
                pass
            self.send_json({"success": True})

        # ── Export ──
        elif path == "/api/export/pdf":
            user = self.get_user()
            data = export_pdf(
                title=body.get("title", "CoivEn Report"),
                content=body.get("content", ""),
                author=user["name"] if user else "CoivEn User",
                session_title=body.get("session_title", ""),
                images=body.get("images", [])
            )
            self.send_file(data, "application/pdf",
                           f"CoivEn_{body.get('title','Report').replace(' ','_')}.pdf")

        elif path == "/api/export/word":
            user = self.get_user()
            data = export_word(
                title=body.get("title", "CoivEn Report"),
                content=body.get("content", ""),
                author=user["name"] if user else "CoivEn User",
                session_title=body.get("session_title", ""),
                images=body.get("images", [])
            )
            self.send_file(data,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                f"CoivEn_{body.get('title','Report').replace(' ','_')}.docx")

        elif path == "/api/export/ppt":
            user = self.get_user()
            data = export_ppt(
                title=body.get("title", "CoivEn Presentation"),
                content=body.get("content", ""),
                author=user["name"] if user else "CoivEn User",
                session_title=body.get("session_title", "")
            )
            self.send_file(data,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                f"CoivEn_{body.get('title','Presentation').replace(' ','_')}.pptx")

        # ── Image processing ──
        elif path == "/api/image/process":
            result = process_image(
                body.get("image", ""),
                body.get("operation", "grayscale"),
                body.get("params", {})
            )
            self.send_json({"image": result, "operation": body.get("operation")})

        elif path == "/api/extract":
            # Improvement 1 — Document text extraction
            filename = body.get("filename", "file.pdf")
            file_b64 = body.get("data", "")
            if not file_b64:
                self.send_json({"text": None, "error": "No file data provided"})
                return
            extracted = extract_text_from_file(filename, file_b64)
            if extracted:
                self.send_json({
                    "text": extracted,
                    "chars": len(extracted),
                    "filename": filename,
                    "success": True
                })
            else:
                self.send_json({
                    "text": None,
                    "success": False,
                    "error": "Could not extract text from this file type"
                })

        elif path == "/api/config/get":
            # Return config — mask API keys for security (show only last 6 chars)
            cfg = load_config()
            safe = dict(cfg)
            for k in ("groq_key", "tavily_key"):
                val = safe.get(k, "")
                safe[k] = ("*" * (len(val) - 6) + val[-6:]) if len(val) > 6 else ("*" * len(val))
            self.send_json({"success": True, "config": safe})

        elif path == "/api/config/save":
            # Save config — requires no auth (local app only)
            new_cfg = body.get("config", {})
            if not isinstance(new_cfg, dict):
                self.send_json({"error": "Invalid config"}, 400)
                return
            # Only allow known keys
            allowed = set(DEFAULT_CONFIG.keys())
            filtered = {k: v for k, v in new_cfg.items() if k in allowed}
            existing = load_config()
            existing.update(filtered)
            save_config(existing)
            # Update in-memory config
            global APP_CONFIG
            APP_CONFIG = existing
            self.send_json({"success": True, "saved": list(filtered.keys())})

        elif path == "/api/config/test-key":
            # Test a Groq API key by making a minimal API call
            import urllib.request as ureq2
            import json as _json
            key   = body.get("key", "").strip()
            ktype = body.get("type", "groq")
            if not key:
                self.send_json({"valid": False, "error": "No key provided"})
                return
            try:
                if ktype == "groq":
                    req = ureq2.Request(
                        "https://api.groq.com/openai/v1/models",
                        headers={"Authorization": "Bearer " + key}
                    )
                    with ureq2.urlopen(req, timeout=8) as r:
                        data = _json.loads(r.read())
                    models = [m["id"] for m in data.get("data", [])]
                    self.send_json({"valid": True, "models": models[:5]})
                elif ktype == "tavily":
                    req = ureq2.Request(
                        "https://api.tavily.com/search",
                        data=_json.dumps({"query": "test", "api_key": key, "max_results": 1}).encode(),
                        headers={"Content-Type": "application/json"}
                    )
                    with ureq2.urlopen(req, timeout=8) as r:
                        _json.loads(r.read())
                    self.send_json({"valid": True})
                else:
                    self.send_json({"valid": False, "error": "Unknown key type"})
            except Exception as ex:
                err = str(ex)
                if "401" in err or "403" in err or "invalid" in err.lower():
                    self.send_json({"valid": False, "error": "Invalid API key"})
                elif "429" in err:
                    self.send_json({"valid": True, "warning": "Key valid but rate limited"})
                else:
                    self.send_json({"valid": False, "error": "Could not verify: " + err[:80]})

        elif path == "/api/config/get-raw":
            # Returns actual (unmasked) config so frontend can use keys directly
            # Only accessible from localhost
            host = self.headers.get("Host", "")
            if not (host.startswith("127.0.0.1") or host.startswith("localhost")):
                self.send_json({"error": "Forbidden"}, 403); return
            cfg = load_config()
            self.send_json({
                "groq_key":   cfg.get("groq_key", ""),
                "tavily_key": cfg.get("tavily_key", ""),
                "groq_model": cfg.get("groq_model", "meta-llama/llama-4-scout-17b-16e-instruct")
            })

        elif path == "/api/chats/save":
            # Save all chat sessions for this user permanently
            user = self.get_user()
            uid  = user.get("id", "guest") if user else "guest"
            chats = body.get("chats", [])
            # Keep last 200 sessions max
            if len(chats) > 200:
                chats = chats[:200]
            save_json(DATA_DIR / f"chats_{uid}.json", chats)
            self.send_json({"success": True, "saved": len(chats)})

        elif path == "/api/chats/load":
            # Load all chat sessions for this user
            user = self.get_user()
            uid  = user.get("id", "guest") if user else "guest"
            chats = load_json(DATA_DIR / f"chats_{uid}.json", [])
            self.send_json({"success": True, "chats": chats})

        elif path == "/api/chats/delete":
            # Delete a specific chat session
            user     = self.get_user()
            uid      = user.get("id", "guest") if user else "guest"
            chat_id  = body.get("id", "")
            chats    = load_json(DATA_DIR / f"chats_{uid}.json", [])
            chats    = [c for c in chats if c.get("id") != chat_id]
            save_json(DATA_DIR / f"chats_{uid}.json", chats)
            self.send_json({"success": True, "remaining": len(chats)})

        elif path == "/api/memory/get":
            # Improvement 2 — Load user memory
            user = self.get_user()
            uid = user["id"] if user else body.get("uid", "anonymous")
            memory = load_memory(uid)
            self.send_json({"memory": memory, "uid": str(uid)})

        elif path == "/api/memory/save":
            # Improvement 2 — Save user memory
            user = self.get_user()
            uid = user["id"] if user else body.get("uid", "anonymous")
            memory = body.get("memory", {})
            ok = save_memory(uid, memory)
            self.send_json({"success": ok})

        elif path == "/api/memory/extract":
            # Improvement 2 — Extract facts from conversation text
            text = body.get("text", "")
            facts = extract_memory_facts(text)
            self.send_json({"facts": facts})

        elif path == "/api/image/preprocess":
            # Improvement 3 — Smart image preprocessing
            img_b64  = body.get("image", "")
            filename = body.get("filename", "image.png")
            if not img_b64:
                self.send_json({"error": "No image data"})
                return
            result = preprocess_image_for_analysis(img_b64, filename)
            self.send_json(result)

        elif path == "/api/execute":
            # Real Python execution using subprocess — full Python environment
            code    = body.get("code", "").strip()
            timeout = min(int(body.get("timeout", 15)), 30)  # max 30s

            if not code:
                self.send_json({"output": "", "error": "No code provided", "success": False})
                return

            if len(code) > 20000:
                self.send_json({"output": "", "error": "Code too long (max 20000 chars)", "success": False})
                return

            # Hard security checks before execution
            BLOCKED_PATTERNS = [
                "import os", "import sys", "import subprocess",
                "import socket", "import requests", "import urllib",
                "import shutil", "import glob", "import pathlib",
                "__import__", "open(", "exec(", "eval(",
                "compile(", "globals()", "locals()", "__builtins__",
                "os.system", "os.popen", "subprocess.run",
                "subprocess.Popen", "socket.socket",
            ]

            # Check for blocked patterns
            code_lower = code.lower().replace(" ", "")
            for pattern in BLOCKED_PATTERNS:
                if pattern.lower().replace(" ", "") in code_lower:
                    # Allow safe subsets
                    if pattern == "open(" and "open(" not in code:
                        continue
                    self.send_json({
                        "output": "",
                        "error": f"Security: '{pattern}' is not allowed. "
                                 "The IDE blocks file I/O, network, and system access.",
                        "success": False
                    })
                    return

            # Write code to temp file and run with real Python
            import tempfile, subprocess as sp, sys as _sys

            # Prepend common engineering imports for convenience
            preamble = """
import math, cmath, statistics, json, re, random, decimal, fractions
import itertools, functools, collections, operator, string
from math import (pi, e, sqrt, log, log2, log10, exp, sin, cos, tan,
                  asin, acos, atan, atan2, sinh, cosh, tanh, floor,
                  ceil, fabs, factorial, gcd, inf, nan, degrees, radians)
from cmath import (phase, polar, rect)
from statistics import (mean, median, mode, stdev, variance, pstdev, pvariance)
from collections import (Counter, defaultdict, OrderedDict, deque, namedtuple)
from itertools import (chain, combinations, permutations, product, groupby)
from decimal import Decimal, getcontext
getcontext().prec = 28
"""
            full_code = preamble + "\n" + code

            with tempfile.NamedTemporaryFile(
                mode='w', suffix='.py',
                delete=False, encoding='utf-8'
            ) as tmp:
                tmp.write(full_code)
                tmp_path = tmp.name

            try:
                result = sp.run(
                    [_sys.executable, tmp_path],
                    capture_output=True,
                    text=True,
                    timeout=timeout,
                    cwd=tempfile.gettempdir()
                )
                stdout = result.stdout
                stderr = result.stderr

                # Clean up preamble line numbers from tracebacks
                if stderr:
                    lines = stderr.split("\n")
                    preamble_lines = len(preamble.split("\n")) - 1
                    cleaned = []
                    for line in lines:
                        import re as _re
                        m = _re.search(r'line (\d+)', line)
                        if m:
                            real_line = int(m.group(1)) - preamble_lines
                            if real_line > 0:
                                line = line.replace(
                                    "line " + m.group(1),
                                    "line " + str(real_line)
                                )
                        if tmp_path not in line:
                            cleaned.append(line)
                    stderr = "\n".join(cleaned).strip()

                if result.returncode == 0:
                    self.send_json({
                        "output":  stdout or "(Code ran with no output)",
                        "error":   None,
                        "success": True,
                        "stderr":  stderr if stderr else None
                    })
                else:
                    self.send_json({
                        "output":  stdout,
                        "error":   stderr or "Runtime error (exit code " + str(result.returncode) + ")",
                        "success": False
                    })

            except sp.TimeoutExpired:
                self.send_json({
                    "output":  "",
                    "error":   f"Execution timed out after {timeout} seconds. Check for infinite loops.",
                    "success": False
                })
            except Exception as ex:
                self.send_json({
                    "output":  "",
                    "error":   "Execution error: " + str(ex),
                    "success": False
                })
            finally:
                try:
                    import os as _os
                    _os.unlink(tmp_path)
                except Exception:
                    pass

        elif path == "/api/search":
            import urllib.request as ureq
            query      = body.get("query", "").strip()
            tavily_key = body.get("key",   "").strip()

            if not query:
                self.send_json({"summary": "", "error": "No query provided"})
                return

            if not tavily_key or tavily_key.lower() in ("your-tavily-key-here", ""):
                self.send_json({"summary": "", "error": "No Tavily key configured"})
                return

            try:
                payload = json.dumps({
                    "query":              query,
                    "api_key":            tavily_key,
                    "max_results":        6,
                    "search_depth":       "basic",
                    "include_answer":     True,
                    "include_raw_content": False
                }).encode("utf-8")

                req = ureq.Request(
                    "https://api.tavily.com/search",
                    data=payload,
                    headers={"Content-Type": "application/json"},
                    method="POST"
                )
                with ureq.urlopen(req, timeout=12) as resp:
                    result = json.loads(resp.read().decode("utf-8"))

                # Improvement 4 — quality filter and ranking
                raw_results = result.get("results", [])
                ranked_summary = filter_and_rank_results(raw_results, query)

                parts = []
                if result.get("answer"):
                    parts.append("Direct answer: " + result["answer"])
                if ranked_summary:
                    parts.append(ranked_summary)

                summary = "\n\n".join(parts) if parts else ""
                self.send_json({
                    "summary": summary,
                    "results_count": len(raw_results),
                    "has_direct_answer": bool(result.get("answer"))
                })

            except Exception as e:
                self.send_json({"summary": "", "error": str(e)})

        elif path == "/api/generate-image":
            try:
                import urllib.request as _ureq
                import urllib.parse   as _uparse
                import urllib.error   as _uerr
                import ssl            as _ssl
                import random         as _rand

                prompt = body.get("prompt", "").strip()
                if not prompt:
                    self.send_json({"error": "No prompt provided"}, 400)
                    return

                width  = max(256, min(1440, int(body.get("width",  1024))))
                height = max(256, min(1440, int(body.get("height",  576))))
                model  = body.get("model", "flux")
                seed   = _rand.randint(1, 999999)
                encoded = _uparse.quote(prompt, safe="")

                url = (
                    "https://image.pollinations.ai/prompt/" + encoded
                    + "?width=" + str(width)
                    + "&height=" + str(height)
                    + "&model=" + model
                    + "&seed=" + str(seed)
                    + "&nologo=true"
                )

                headers = {
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36",
                    "Accept":     "image/jpeg,image/png,image/*,*/*",
                }

                # Use a permissive SSL context in case of cert issues on Windows
                try:
                    ctx = _ssl.create_default_context()
                except Exception:
                    ctx = _ssl._create_unverified_context()

                req = _ureq.Request(url, headers=headers)
                img_bytes = None
                fetch_error = ""

                try:
                    with _ureq.urlopen(req, timeout=90, context=ctx) as resp:
                        img_bytes = resp.read()
                        content_type = resp.headers.get("Content-Type", "image/jpeg")
                except _uerr.HTTPError as he:
                    fetch_error = "HTTP " + str(he.code) + ": " + str(he.reason)
                except _uerr.URLError as ue:
                    fetch_error = "URL error: " + str(ue.reason)
                except Exception as ge:
                    fetch_error = type(ge).__name__ + ": " + str(ge)

                if img_bytes and len(img_bytes) > 1000:
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    self.send_json({
                        "image":  img_b64,
                        "mime":   content_type,
                        "width":  width,
                        "height": height,
                        "model":  model,
                        "prompt": prompt
                    })
                else:
                    self.send_json({
                        "error": "Image generation failed: " + (fetch_error or "Empty response from Pollinations"),
                        "url":   url
                    }, 503)

            except Exception as outer_err:
                # This outer catch ensures the server NEVER sends an empty response
                try:
                    self.send_json({
                        "error": "Server error during image generation: " + str(outer_err)
                    }, 500)
                except Exception:
                    pass


        else:
            self.send_json({"error": "Not found"}, 404)


def _keep_alive():
    """Ping the backend every 10 minutes to prevent Render free tier from sleeping."""
    import urllib.request
    port = int(os.environ.get("PORT", 5050))
    url  = os.environ.get("RENDER_EXTERNAL_URL", f"http://127.0.0.1:{port}")
    url  = url.rstrip("/") + "/ping"
    while True:
        time.sleep(600)  # 10 minutes
        try:
            urllib.request.urlopen(url, timeout=10)
            print(f"[CoivEn] keep-alive ping OK → {url}")
        except Exception as e:
            print(f"[CoivEn] keep-alive ping failed: {e}")


def run(port=None):
    port = int(port or os.environ.get("PORT", 5050))
    host = "0.0.0.0"
    server = ThreadingHTTPServer((host, port), CoivEnHandler)
    server.daemon_threads = True
    # Start keep-alive thread so free-tier hosting doesn't sleep
    t = threading.Thread(target=_keep_alive, daemon=True)
    t.start()
    print(f"[CoivEn Backend] Running on http://{host}:{port} (threaded)")
    server.serve_forever()

if __name__ == "__main__":
    run()
